import os
import json
import docx
import openpyxl
from pptx import Presentation
from PIL import Image
import pdfplumber
from paddleocr import PaddleOCR

# 初始化 PaddleOCR (中+英，啟用文字方向檢測)
ocr = PaddleOCR(use_textline_orientation=True, lang='ch')

def ocr_image(img_path_or_pil):
    """使用 PaddleOCR 辨識圖片，回傳文字內容"""
    if isinstance(img_path_or_pil, str):
        results = ocr.ocr(img_path_or_pil)
    else:
        temp_path = "temp_ocr_image.png"
        img_path_or_pil.save(temp_path)
        results = ocr.ocr(temp_path)
        os.remove(temp_path)

    lines = []
    if results and results[0]:
        for line in results[0]:
            text = line[1][0]
            lines.append(text)
    return "\n".join(lines) if lines else ""

def read_file_content(filepath):
    file_extension = os.path.splitext(filepath)[1].lower()
    content = ""

    try:
        if file_extension == '.docx':
            doc = docx.Document(filepath)
            for paragraph in doc.paragraphs:
                content += paragraph.text + '\n'
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    content += '\t'.join(row_text) + '\n'

        elif file_extension == '.pdf':
            pdf_content_parts = []
            with pdfplumber.open(filepath) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        pdf_content_parts.append(text)
                    else:
                        page_image = page.to_image(resolution=300).original
                        ocr_text = ocr_image(page_image)
                        if ocr_text.strip():
                            pdf_content_parts.append(f"[第 {i+1} 頁 OCR]\n{ocr_text.strip()}")
                        else:
                            pdf_content_parts.append(f"[第 {i+1} 頁] 未能提取內容")
            content = "\n---\n".join(pdf_content_parts)

        elif file_extension in ['.xlsx', '.xls']:
            workbook = openpyxl.load_workbook(filepath)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            content += str(cell.value) + ' '
                content += '\n'

        elif file_extension in ['.pptx', '.ppt']:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + '\n'

        elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            ocr_text = ocr_image(filepath)
            if ocr_text.strip():
                content = f"圖片檔案 (OCR):\n{ocr_text.strip()}"
            else:
                img = Image.open(filepath)
                content = f"圖片檔案, 格式: {img.format}, 大小: {img.size}, 模式: {img.mode} (OCR 無內容)"

        elif file_extension == '.txt':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()

        elif file_extension in ['.zip', '.rar', '.7z', '.tar', '.gz']:
            content = f"壓縮檔案, 無法直接提取: {file_extension}"

        else:
            content = f"不支援的檔案類型: {file_extension}"

    except Exception as e:
        content = f"讀取檔案時發生錯誤: {e}"

    return content

def main():
    if not os.path.exists('db'):
        os.makedirs('db')

    data_to_save = []

    for root, dirs, files in os.walk('data'):
        if root == 'data':
            project_name = 'none'
        else:
            relative_path = os.path.relpath(root, 'data')
            project_name = relative_path.split(os.sep)[0]

        for filename in files:
            if not filename.startswith('.'):
                filepath = os.path.join(root, filename)
                content = read_file_content(filepath)
                entry = {
                    "project": project_name,
                    "filename": filename,
                    "content": content
                }
                data_to_save.append(entry)

    with open('db/O.json', 'w', encoding='utf-8') as f:
        json.dump(data_to_save, f, ensure_ascii=False, indent=4)

    print(f"已儲存至 db/O.json，共處理 {len(data_to_save)} 個檔案。")

if __name__ == "__main__":
    main()
